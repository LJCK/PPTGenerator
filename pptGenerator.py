import sys
sys.path.insert(0,'site-packages')

from pptx import Presentation
import json
import os
import qrcode
import dotenv

user_input_choices = "1. Check the existing URL\n2. Change the URL in the .env file.\n3. Generate a ppt based on an existing URL stored in the .env file.\n4. End\nYour choice: "
dotenv_file = dotenv.find_dotenv()
dotenv.load_dotenv(dotenv_file)


def generate_ppt(URL, FIELD_ID):

  course_name = input("\nEnter the course name:\n")
  special_characters = '\/:*?|"<>'
  for z in special_characters:
    if z in course_name:
      print("Course name cannot contain any of the follower characters: ", special_characters)
      print("File is not saved, returning to main page.")
      return

  if course_name.isspace():
    print("Your course name cannot only contains whitespace.")
    print("File is not saved, returning to main page.")
    return

  new_course_name = "%20".join(course_name.split())
  string_URL = URL+"?"+FIELD_ID+"="+new_course_name
  
  prs = Presentation("template.pptx")
  slide = prs.slides[0]

  shapes =[shape for shape in slide.shapes]
  string = old_QR_code = None

  for i in shapes:
    if i.name == "URL":
      string = i
    elif i.name == "QRCode":
      old_QR_code = i

  if string == None:
    print("No text box for URL found. Returning to main page.\n")
    return 
  if old_QR_code == None:
    print("No image box for QR code found. Returning to main page.\n")
    return

  cur_text = string.text_frame.paragraphs[0].runs[0].text
  new_text = cur_text.replace(cur_text,string_URL)
  string.text_frame.paragraphs[0].runs[0].text = new_text
  
  img = qrcode.make(string_URL)
  img_file_name = course_name + ".png"
  img.save(img_file_name)

  with open(img_file_name, 'rb') as f:
    rImgBlob = f.read()

  imgPic = old_QR_code._pic
  imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
  imgPart = slide.part.related_part(imgRID)
  imgPart._blob = rImgBlob

  os.remove(img_file_name)
  prs.save(course_name + ".pptx")
  print(course_name, ".pptx is saved to the following path: ", os.path.abspath(os.getcwd()), "\n")
    
def change_prefix_URL(URL, FIELD_ID):
  
    print("\nThe existing URL: ", URL, " field ID: ", FIELD_ID)
    user_input = input("Choose the one you want to change \n1. URL\n2. field ID. \n3. return \n")

    if user_input == "1":
      # change URL
      newURL = input("New URL: ")
      os.environ["URL"] = newURL
      dotenv.set_key(dotenv_file, "URL", os.environ["URL"])      
      print("New URL updated. Returning to main page.\n")

    elif user_input == "2":
      # change field id
      newField_id = input("New field id: ")
      os.environ["field_id"] = newField_id
      dotenv.set_key(dotenv_file, "field_id", os.environ["field_id"])
      print("New field id updated. Returning to main page.\n")

    elif user_input == '3':
      return

    else:
      print("Wrong input. Returning to main page\n")
      return

if __name__ == "__main__":
  while(1):
    print("#################################################################")
    URL = os.environ["URL"]
    FIELD_ID = os.environ["FIELD_ID"]
    user_input = input(user_input_choices)

    if user_input == '1':
      print("\nThe existing URL: ", URL, " field id: ", FIELD_ID, "\n")
      
    elif user_input == '2':
      # change URL pre-fix
      if URL and FIELD_ID:
        change_prefix_URL(URL, FIELD_ID)
      else:
        print("\nThe .env file does not contain URL or field_id, please check!\n")

    elif user_input == '3':
      if URL and FIELD_ID:
        generate_ppt(URL, FIELD_ID)
      else:
        print("\nThe .env file does not contain URL or field_id, please check!\n")

    elif user_input == '4':
      break

    else:
      print("\nYou can only choose one from 1-4. \n")
      